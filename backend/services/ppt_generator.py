import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from typing import Dict, Any

# ── Color constants ────────────────────────────────────────────────────────────
DARK_NAVY      = RGBColor(15,  23,  42)   # #0F172A
NAVY_BLUE      = RGBColor(30,  58,  138)  # #1E3A8A
MID_NAVY       = RGBColor(30,  64,  175)  # #1E40AF
LIGHT_BG       = RGBColor(248, 250, 252)  # #F8FAFC
GRAY_BG        = RGBColor(226, 232, 240)  # #E2E8F0
ALT_ROW        = RGBColor(241, 245, 249)  # #F1F5F9
WHITE          = RGBColor(255, 255, 255)
RED_COLOR      = RGBColor(220,  38,  38)  # #DC2626
GREEN_COLOR    = RGBColor(22,  163,  74)  # #16A34A
ORANGE_COLOR   = RGBColor(217, 119,   6)  # #D97706
BLUE_COLOR     = RGBColor(37,   99, 235)  # #2563EB
INDIGO_COLOR   = RGBColor(79,   70, 229)  # #4F46E5
AMBER_COLOR    = RGBColor(245, 158,  11)  # #F59E0B
PINK_HIGHLIGHT = RGBColor(254, 226, 226)  # #FEE2E2
BORDER_GRAY    = RGBColor(203, 213, 225)  # #CBD5E1
INSIGHT_BG     = RGBColor(239, 246, 255)  # #EFF6FF — soft blue insight box


# ── Helpers ────────────────────────────────────────────────────────────────────

def _set_cell_text(cell, text, bold=False, color=DARK_NAVY, size_pt=10,
                   align=PP_ALIGN.LEFT, italic=False, wrap=True):
    tf = cell.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.text = str(text)
    p.alignment = align
    run = p.runs[0] if p.runs else p.add_run()
    run.font.size  = Pt(size_pt)
    run.font.bold  = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name  = "Calibri"


def _cell_fill(cell, color: RGBColor):
    cell.fill.solid()
    cell.fill.fore_color.rgb = color


def _get_safe_logo_path(logo_path: str = None) -> str:
    """Ensure logo is in a python-pptx compatible format (PNG/JPEG)."""
    if not logo_path or not os.path.exists(logo_path):
        # Check standard PNG fallback
        default_png = os.path.join("frontend", "public", "AMC Logo.png")
        if os.path.exists(default_png):
            return default_png
        return None

    if logo_path.lower().endswith(".webp"):
        png_path = logo_path[:-5] + ".png"
        if not os.path.exists(png_path):
            try:
                from PIL import Image
                im = Image.open(logo_path)
                im.save(png_path, "PNG")
            except Exception:
                return None
        return png_path
    return logo_path


def _add_header(slide, title_text: str, subtitle_text: str, logo_path: str = None):
    """Full-width dark navy header bar with optional AMC logo."""
    W = Inches(13.333)
    H = Inches(1.1)
    banner = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    Inches(0), Inches(0), W, H)
    banner.fill.solid()
    banner.fill.fore_color.rgb = DARK_NAVY
    banner.line.color.rgb      = DARK_NAVY

    tf = banner.text_frame
    tf.word_wrap     = True
    tf.margin_left   = Inches(0.45)
    tf.margin_top    = Inches(0.12)
    tf.margin_right  = Inches(1.2)

    p1 = tf.paragraphs[0]
    p1.text = title_text
    if p1.runs:
        p1.runs[0].font.name  = "Calibri"
        p1.runs[0].font.size  = Pt(20)
        p1.runs[0].font.bold  = True
        p1.runs[0].font.color.rgb = WHITE

    p2 = tf.add_paragraph()
    p2.text = subtitle_text
    if p2.runs:
        p2.runs[0].font.name  = "Calibri"
        p2.runs[0].font.size  = Pt(10)
        p2.runs[0].font.color.rgb = RGBColor(148, 163, 184)

    # Optional logo safely loaded
    safe_logo = _get_safe_logo_path(logo_path)
    if safe_logo and os.path.exists(safe_logo):
        slide.shapes.add_picture(safe_logo,
                                 Inches(12.35), Inches(0.12),
                                 Inches(0.85), Inches(0.85))


def _add_kpi_card(slide, left, top, width, height,
                  label, value, sub="", accent=BLUE_COLOR):
    """Rounded KPI card with accent top-bar."""
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = LIGHT_BG
    card.line.color.rgb      = BORDER_GRAY

    # Accent stripe
    stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    left, top, width, Inches(0.08))
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = accent
    stripe.line.color.rgb      = accent

    tf = card.text_frame
    tf.word_wrap    = True
    tf.margin_top   = Inches(0.18)
    tf.margin_left  = Inches(0.18)
    tf.margin_right = Inches(0.1)

    p0 = tf.paragraphs[0]
    p0.text = label.upper()
    r0 = p0.runs[0] if p0.runs else p0.add_run()
    r0.font.size  = Pt(9)
    r0.font.bold  = True
    r0.font.color.rgb = RGBColor(100, 116, 139)
    r0.font.name  = "Calibri"

    p1 = tf.add_paragraph()
    p1.text = str(value)
    r1 = p1.runs[0] if p1.runs else p1.add_run()
    r1.font.size  = Pt(22)
    r1.font.bold  = True
    r1.font.color.rgb = DARK_NAVY
    r1.font.name  = "Calibri"

    if sub:
        p2 = tf.add_paragraph()
        p2.text = str(sub)
        r2 = p2.runs[0] if p2.runs else p2.add_run()
        r2.font.size  = Pt(9)
        r2.font.bold  = True
        r2.font.color.rgb = accent
        r2.font.name  = "Calibri"


def _add_insights_box(slide, left, top, width, height,
                      title_text: str, bullets: list):
    """Styled insights box with indigo left border + bullet points."""
    # Background
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = INSIGHT_BG
    box.line.color.rgb      = BORDER_GRAY

    # Left accent strip
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    left, top, Inches(0.08), height)
    accent.fill.solid()
    accent.fill.fore_color.rgb = INDIGO_COLOR
    accent.line.color.rgb      = INDIGO_COLOR

    tf = box.text_frame
    tf.word_wrap    = True
    tf.margin_left  = Inches(0.28)
    tf.margin_top   = Inches(0.12)
    tf.margin_right = Inches(0.15)

    p0 = tf.paragraphs[0]
    p0.text = title_text
    r0 = p0.runs[0] if p0.runs else p0.add_run()
    r0.font.name  = "Calibri"
    r0.font.size  = Pt(11)
    r0.font.bold  = True
    r0.font.color.rgb = NAVY_BLUE

    for bullet in bullets:
        pb = tf.add_paragraph()
        pb.text = f"\u2022  {bullet}"
        pb.space_before = Pt(4)
        rb = pb.runs[0] if pb.runs else pb.add_run()
        rb.font.name  = "Calibri"
        rb.font.size  = Pt(9)
        rb.font.color.rgb = DARK_NAVY


def _style_table_header(table, col_count, labels, col_widths=None):
    """Apply dark navy header row to table."""
    for c, h in enumerate(labels):
        cell = table.cell(0, c)
        _cell_fill(cell, NAVY_BLUE)
        _set_cell_text(cell, h, bold=True, color=WHITE, size_pt=9,
                       align=PP_ALIGN.CENTER)
    if col_widths:
        for c, w in enumerate(col_widths):
            table.columns[c].width = w


# ──────────────────────────────────────────────────────────────────────────────
# MAIN BUILD FUNCTION
# ──────────────────────────────────────────────────────────────────────────────

def build_ppt_presentation(
    road_stats:     Dict[str, Any],
    drainage_stats: Dict[str, Any],
    water_stats:    Dict[str, Any],
    chart_paths:    Dict[str, Any],
    date_range:     str,
    output_ppt_path: str,
    logo_path:      str = None
):
    prs = Presentation()
    prs.slide_width  = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 1 — Road: Chart (top) + Table start (bottom)
    # ══════════════════════════════════════════════════════════════════════════
    s1 = prs.slides.add_slide(blank)
    _add_header(s1,
                "Zone wise Road CCRS Complaints Report",
                f"Open Complaints Breakdown by Problem Category across Zones  \u2022  Reporting Period: {date_range}",
                logo_path)

    # Chart — left half, right has KPI cards
    chart_top = Inches(1.2)
    if "road_chart" in chart_paths and os.path.exists(chart_paths["road_chart"]):
        s1.shapes.add_picture(chart_paths["road_chart"],
                              Inches(0.3), chart_top, Inches(8.6), Inches(4.1))

    # KPI cards — right column
    _add_kpi_card(s1, Inches(9.2), Inches(1.3), Inches(3.8), Inches(1.15),
                  "TOTAL COMPLAINTS", road_stats["grand_total"],
                  f"{road_stats['pct_open']}% Open Rate", RED_COLOR)
    _add_kpi_card(s1, Inches(9.2), Inches(2.6), Inches(3.8), Inches(1.15),
                  "CLOSED COMPLAINTS", road_stats["total_closed"],
                  f"{road_stats['resolution_rate']}% Resolution Rate", GREEN_COLOR)
    _add_kpi_card(s1, Inches(9.2), Inches(3.9), Inches(3.8), Inches(1.15),
                  "OPEN COMPLAINTS", road_stats["total_open"],
                  "Requires Field Maintenance", ORANGE_COLOR)

    # Road Insights box
    _add_insights_box(s1, Inches(0.3), Inches(5.5), Inches(12.733), Inches(1.75),
                      "KEY INSIGHTS & EXECUTIVE SUMMARY",
                      road_stats.get("insights", []))

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 2 — Road: Full Data Table
    # ══════════════════════════════════════════════════════════════════════════
    s2 = prs.slides.add_slide(blank)
    _add_header(s2,
                "Zone wise Road CCRS Complaints - Detailed Data Table",
                f"Zone & Problem Category Breakdown  \u2022  {date_range}",
                logo_path)

    zones = road_stats["zones"]
    # Count rows: each category row + 1 subtotal row per zone + 1 header + 1 grand total
    n_rows = sum(len(z["categories"]) + 1 for z in zones) + 2
    tbl2 = s2.shapes.add_table(n_rows, 5,
                                Inches(0.3), Inches(1.25),
                                Inches(12.733), Inches(5.95)).table
    col_widths = [Inches(2.3), Inches(4.1), Inches(2.0), Inches(2.0), Inches(2.333)]
    for c, w in enumerate(col_widths):
        tbl2.columns[c].width = w

    _style_table_header(tbl2, 5,
                        ["ZONE", "PROBLEM CATEGORY", "CLOSED", "OPEN", "GRAND TOTAL"])

    r = 1
    for z_idx, z in enumerate(zones):
        row_bg = ALT_ROW if z_idx % 2 == 0 else WHITE
        for cat_name, cat_vals in z["categories"].items():
            _cell_fill(tbl2.cell(r, 0), row_bg)
            _cell_fill(tbl2.cell(r, 1), row_bg)
            _cell_fill(tbl2.cell(r, 2), row_bg)
            _cell_fill(tbl2.cell(r, 3), row_bg)
            _cell_fill(tbl2.cell(r, 4), row_bg)

            _set_cell_text(tbl2.cell(r, 0), z["zone"], size_pt=9)
            _set_cell_text(tbl2.cell(r, 1), cat_name, size_pt=9)
            _set_cell_text(tbl2.cell(r, 2), cat_vals["closed"],
                           size_pt=9, align=PP_ALIGN.CENTER)
            _set_cell_text(tbl2.cell(r, 3), cat_vals["open"],
                           bold=True, color=RED_COLOR, size_pt=9, align=PP_ALIGN.CENTER)
            _set_cell_text(tbl2.cell(r, 4), cat_vals["grand_total"],
                           size_pt=9, align=PP_ALIGN.CENTER)
            r += 1

        # Subtotal row
        for c in range(5):
            _cell_fill(tbl2.cell(r, c), GRAY_BG)
        _set_cell_text(tbl2.cell(r, 0), f"{z['zone']} — Subtotal",
                       bold=True, size_pt=9)
        _set_cell_text(tbl2.cell(r, 1), "", size_pt=9)
        _set_cell_text(tbl2.cell(r, 2), z["subtotal_closed"],
                       bold=True, size_pt=9, align=PP_ALIGN.CENTER)
        _set_cell_text(tbl2.cell(r, 3), z["subtotal_open"],
                       bold=True, color=RED_COLOR, size_pt=9, align=PP_ALIGN.CENTER)
        _set_cell_text(tbl2.cell(r, 4), z["subtotal_grand_total"],
                       bold=True, size_pt=9, align=PP_ALIGN.CENTER)
        r += 1

    # Grand Total row
    for c in range(5):
        _cell_fill(tbl2.cell(r, c), DARK_NAVY)
    _set_cell_text(tbl2.cell(r, 0), "GRAND TOTAL",
                   bold=True, color=WHITE, size_pt=10)
    _set_cell_text(tbl2.cell(r, 1), "", color=WHITE, size_pt=9)
    _set_cell_text(tbl2.cell(r, 2), road_stats["total_closed"],
                   bold=True, color=WHITE, size_pt=10, align=PP_ALIGN.CENTER)
    _set_cell_text(tbl2.cell(r, 3), road_stats["total_open"],
                   bold=True, color=RGBColor(252, 165, 165), size_pt=10, align=PP_ALIGN.CENTER)
    _set_cell_text(tbl2.cell(r, 4), road_stats["grand_total"],
                   bold=True, color=WHITE, size_pt=10, align=PP_ALIGN.CENTER)

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 3 — Drainage: KPIs + Two Charts
    # ══════════════════════════════════════════════════════════════════════════
    s3 = prs.slides.add_slide(blank)
    _add_header(s3,
                "CCRS Drainage Complaints Summary Report",
                f"Comprehensive Zone-Wise Open Complaints Analysis  \u2022  Period: {date_range}",
                logo_path)

    card_w = Inches(3.0)
    card_h = Inches(1.1)
    card_y = Inches(1.2)
    top_z  = drainage_stats["highest_open_zone"]
    _add_kpi_card(s3, Inches(0.3),  card_y, card_w, card_h,
                  "TOTAL COMPLAINTS", drainage_stats["grand_total"], accent=BLUE_COLOR)
    _add_kpi_card(s3, Inches(3.55), card_y, card_w, card_h,
                  "CLOSED COMPLAINTS", drainage_stats["total_closed"],
                  f"{drainage_stats['resolution_rate']}% Res. Rate", GREEN_COLOR)
    _add_kpi_card(s3, Inches(6.8),  card_y, card_w, card_h,
                  "OPEN COMPLAINTS", drainage_stats["total_open"],
                  f"{drainage_stats['pct_open']}% Pending", RED_COLOR)
    _add_kpi_card(s3, Inches(10.05), card_y, card_w, card_h,
                  "HIGHEST OPEN ZONE", top_z["zone"],
                  f"{top_z['total_open']} Open Cases", ORANGE_COLOR)

    # Two charts side-by-side
    chart_y = Inches(2.45)
    ch = Inches(4.8)
    if "drainage_cat" in chart_paths and os.path.exists(chart_paths["drainage_cat"]):
        s3.shapes.add_picture(chart_paths["drainage_cat"],
                              Inches(0.3), chart_y, Inches(6.3), ch)
    if "drainage_total" in chart_paths and os.path.exists(chart_paths["drainage_total"]):
        s3.shapes.add_picture(chart_paths["drainage_total"],
                              Inches(6.9), chart_y, Inches(6.133), ch)

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 4 — Drainage: Data Table + Insights
    # ══════════════════════════════════════════════════════════════════════════
    s4 = prs.slides.add_slide(blank)
    _add_header(s4,
                "Zone-Wise CCRS Drainage Complaints Detailed Data Table",
                f"Zone Matrix with Category Open Counts & Resolution Rates  \u2022  {date_range}",
                logo_path)

    d_rows = drainage_stats["table_rows"]
    d_cats = drainage_stats["categories"]
    total_cols = len(d_cats) + 5  # Zone + cats + Total Open + Total Closed + Grand Total + % Open

    tbl4 = s4.shapes.add_table(len(d_rows) + 2, total_cols,
                                Inches(0.3), Inches(1.25),
                                Inches(12.733), Inches(3.8)).table

    d_headers = ["Zone"] + d_cats + ["Total Open", "Total Closed", "Grand Total", "% Open"]
    zone_w    = Inches(1.4)
    cat_w     = Inches(9.333 / max(len(d_cats), 1))
    summary_w = Inches(1.2)

    tbl4.columns[0].width = zone_w
    for c in range(1, len(d_cats) + 1):
        tbl4.columns[c].width = cat_w
    for c in range(len(d_cats) + 1, total_cols):
        tbl4.columns[c].width = summary_w

    # Header row
    for c, h in enumerate(d_headers):
        cell = tbl4.cell(0, c)
        _cell_fill(cell, NAVY_BLUE)
        _set_cell_text(cell, h, bold=True, color=WHITE, size_pt=8, align=PP_ALIGN.CENTER)

    # Data rows
    for r_idx, row in enumerate(d_rows):
        rp = r_idx + 1
        row_bg = ALT_ROW if r_idx % 2 == 0 else WHITE

        for c in range(total_cols):
            _cell_fill(tbl4.cell(rp, c), row_bg)

        _set_cell_text(tbl4.cell(rp, 0), row["zone"], bold=True, size_pt=8)
        for c_idx, cat in enumerate(d_cats):
            val = row["cat_open"].get(cat, 0)
            _set_cell_text(tbl4.cell(rp, c_idx + 1), val,
                           size_pt=8, align=PP_ALIGN.CENTER)

        c_off = len(d_cats) + 1
        _set_cell_text(tbl4.cell(rp, c_off),   row["total_open"],
                       bold=True, color=RED_COLOR, size_pt=8, align=PP_ALIGN.CENTER)
        _set_cell_text(tbl4.cell(rp, c_off+1), row["total_closed"],
                       size_pt=8, align=PP_ALIGN.CENTER)
        _set_cell_text(tbl4.cell(rp, c_off+2), row["grand_total"],
                       size_pt=8, align=PP_ALIGN.CENTER)
        _set_cell_text(tbl4.cell(rp, c_off+3), f"{row['pct_open']}%",
                       size_pt=8, align=PP_ALIGN.CENTER)

    # Grand total row
    last4 = len(d_rows) + 1
    for c in range(total_cols):
        _cell_fill(tbl4.cell(last4, c), DARK_NAVY)
    _set_cell_text(tbl4.cell(last4, 0), "ALL ZONES", bold=True, color=WHITE, size_pt=8)
    for c_idx, cat in enumerate(d_cats):
        _set_cell_text(tbl4.cell(last4, c_idx + 1),
                       drainage_stats["cat_totals"].get(cat, 0),
                       bold=True, color=WHITE, size_pt=8, align=PP_ALIGN.CENTER)
    c_off = len(d_cats) + 1
    _set_cell_text(tbl4.cell(last4, c_off),   drainage_stats["total_open"],
                   bold=True, color=RGBColor(252, 165, 165), size_pt=8, align=PP_ALIGN.CENTER)
    _set_cell_text(tbl4.cell(last4, c_off+1), drainage_stats["total_closed"],
                   bold=True, color=WHITE, size_pt=8, align=PP_ALIGN.CENTER)
    _set_cell_text(tbl4.cell(last4, c_off+2), drainage_stats["grand_total"],
                   bold=True, color=WHITE, size_pt=8, align=PP_ALIGN.CENTER)
    _set_cell_text(tbl4.cell(last4, c_off+3), f"{drainage_stats['pct_open']}%",
                   bold=True, color=WHITE, size_pt=8, align=PP_ALIGN.CENTER)

    # Insights box
    _add_insights_box(s4, Inches(0.3), Inches(5.2), Inches(12.733), Inches(2.05),
                      "KEY INSIGHTS & EXECUTIVE SUMMARY",
                      drainage_stats.get("insights", []))

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 5 — Water: KPIs + Matrix Table + Chart
    # ══════════════════════════════════════════════════════════════════════════
    s5 = prs.slides.add_slide(blank)
    _add_header(s5,
                "CCRS Water Complaints: Zone & Category Open Summary",
                f"Reporting Period: {date_range}",
                logo_path)

    w_cats  = water_stats["categories"]
    w_zrows = water_stats["zone_rows"]
    col_totals = water_stats["col_totals"]

    # KPI Cards
    _add_kpi_card(s5, Inches(0.3),  Inches(1.2), Inches(3.0), Inches(1.05),
                  "TOTAL OPEN COMPLAINTS", water_stats["total_open"],
                  "Citywide Water Backlog", RED_COLOR)
    _add_kpi_card(s5, Inches(3.55), Inches(1.2), Inches(3.0), Inches(1.05),
                  "TOP OPEN ZONE", water_stats["top_zone"]["zone"],
                  f"{water_stats['top_zone']['total_open']} Open Cases", ORANGE_COLOR)
    _add_kpi_card(s5, Inches(6.8),  Inches(1.2), Inches(3.0), Inches(1.05),
                  "TOP OPEN CATEGORY", water_stats["top_category"]["name"],
                  f"{water_stats['top_category']['count']} Cases ({water_stats['top_category']['pct']}%)", BLUE_COLOR)
    _add_kpi_card(s5, Inches(10.05), Inches(1.2), Inches(2.978), Inches(1.05),
                  "CATEGORIES TRACKED", len(w_cats),
                  "Active Water Issue Types", GREEN_COLOR)

    # Compact Water Matrix Table
    n_w_cols = len(w_cats) + 2  # Zone + cats + Total Open
    tbl5 = s5.shapes.add_table(len(w_zrows) + 2, n_w_cols,
                                Inches(0.3), Inches(2.38),
                                Inches(12.733), Inches(2.55)).table

    zone_w5 = Inches(1.5)
    cat_w5  = Inches(10.233 / max(len(w_cats), 1))
    tot_w5  = Inches(1.0)
    tbl5.columns[0].width = zone_w5
    for c in range(1, len(w_cats) + 1):
        tbl5.columns[c].width = cat_w5
    tbl5.columns[n_w_cols - 1].width = tot_w5

    # Header
    w_headers = ["Zone Name"] + w_cats + ["Total Open"]
    for c, h in enumerate(w_headers):
        cell = tbl5.cell(0, c)
        _cell_fill(cell, NAVY_BLUE)
        _set_cell_text(cell, h, bold=True, color=WHITE, size_pt=8, align=PP_ALIGN.CENTER)

    # Find column peak values for pink highlight
    col_peaks = {}
    for cat in w_cats:
        vals = [z["categories"].get(cat, 0) for z in w_zrows]
        col_peaks[cat] = max(vals) if vals else 0

    for r_idx, zr in enumerate(w_zrows):
        rp = r_idx + 1
        row_bg = ALT_ROW if r_idx % 2 == 0 else WHITE

        for c in range(n_w_cols):
            _cell_fill(tbl5.cell(rp, c), row_bg)

        _set_cell_text(tbl5.cell(rp, 0), zr["zone"], bold=True, size_pt=8)

        for c_idx, cat in enumerate(w_cats):
            val  = zr["categories"].get(cat, 0)
            cell = tbl5.cell(rp, c_idx + 1)
            is_peak = val > 0 and val == col_peaks.get(cat, -1)
            if is_peak:
                _cell_fill(cell, PINK_HIGHLIGHT)
                _set_cell_text(cell, val, bold=True, color=RED_COLOR,
                               size_pt=8, align=PP_ALIGN.CENTER)
            else:
                _set_cell_text(cell, val, size_pt=8, align=PP_ALIGN.CENTER)

        _set_cell_text(tbl5.cell(rp, n_w_cols - 1), zr["total_open"],
                       bold=True, color=RED_COLOR, size_pt=8, align=PP_ALIGN.CENTER)

    # Total row
    last5 = len(w_zrows) + 1
    for c in range(n_w_cols):
        _cell_fill(tbl5.cell(last5, c), DARK_NAVY)
    _set_cell_text(tbl5.cell(last5, 0), "Total Open", bold=True, color=WHITE, size_pt=8)
    for c_idx, cat in enumerate(w_cats):
        _set_cell_text(tbl5.cell(last5, c_idx + 1),
                       col_totals.get(cat, 0),
                       bold=True, color=WHITE, size_pt=8, align=PP_ALIGN.CENTER)
    _set_cell_text(tbl5.cell(last5, n_w_cols - 1), water_stats["total_open"],
                   bold=True, color=RGBColor(252, 165, 165), size_pt=8, align=PP_ALIGN.CENTER)

    # Water Bar Chart
    if "water_chart" in chart_paths and os.path.exists(chart_paths["water_chart"]):
        s5.shapes.add_picture(chart_paths["water_chart"],
                              Inches(0.3), Inches(5.05), Inches(8.5), Inches(2.2))

    # Insights box (right of chart)
    _add_insights_box(s5, Inches(9.05), Inches(5.05), Inches(4.0), Inches(2.2),
                      "KEY INSIGHTS",
                      water_stats.get("insights", [])[:3])

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 6 — Water: Strategic Focus & Full Insights
    # ══════════════════════════════════════════════════════════════════════════
    s6 = prs.slides.add_slide(blank)
    _add_header(s6,
                "CCRS Water Complaints - Strategic Focus & Executive Summary",
                f"Actionable Insights Based on Zone x Category Analysis  \u2022  {date_range}",
                logo_path)

    _add_insights_box(s6, Inches(0.3), Inches(1.3), Inches(12.733), Inches(5.8),
                      "KEY INSIGHTS & STRATEGIC FOCUS AREAS",
                      water_stats.get("insights", []))

    prs.save(output_ppt_path)
    return output_ppt_path
